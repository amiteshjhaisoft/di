# Author: Amitesh Jha | iSoft | 2025-10-03 
#
# LLM chat with Local RAG (Chroma) + polished ChatGPT-style UI.
# Uses branded avatars:
#   - User:       assets/avatar.png
#   - Assistant:  assets/Forecast360.png
#
# Run:
#   pip install streamlit chromadb sentence-transformers pypdf python-docx pandas openpyxl anthropic requests
#   streamlit run app_sidebar.py
#
# Env (optional):
#   ANTHROPIC_API_KEY         -> for Claude
#   ISOFT_LOGO_PATH           -> path to iSOFT logo
#   USER_AVATAR_PATH          -> override user avatar path
#   ASSISTANT_AVATAR_PATH     -> override assistant avatar path
#   (Ollama runs at http://localhost:11434)

from __future__ import annotations

import os, re, io, gc, glob, time, uuid, hashlib, logging, base64
from dataclasses import dataclass
from typing import List, Dict, Optional, Tuple
from pathlib import Path

import requests
import pandas as pd
import streamlit as st

# -------------------------------------------------------------------
# Paths & helpers
# -------------------------------------------------------------------
try:
    APP_DIR = Path(__file__).resolve().parent
except NameError:
    APP_DIR = Path.cwd()

def _resolve_logo_path() -> Optional[Path]:
    env_logo = os.getenv("ISOFT_LOGO_PATH")
    candidates = [
        APP_DIR / "assets" / "isoft_logo.png",
        Path.cwd() / "assets" / "isoft_logo.png",
        Path("C:/Decision_Intelligence/assets/isoft_logo.png"),
        Path(env_logo).expanduser().resolve() if env_logo else None,
    ]
    for p in candidates:
        if p and p.exists():
            return p
    return None

def _resolve_avatar_paths() -> Tuple[Optional[Path], Optional[Path]]:
    env_user = os.getenv("USER_AVATAR_PATH")
    env_asst = os.getenv("ASSISTANT_AVATAR_PATH")
    user_candidates = [
        APP_DIR / "assets" / "avatar.png",
        Path.cwd() / "assets" / "avatar.png",
        Path(env_user).expanduser().resolve() if env_user else None,
    ]
    asst_candidates = [
        APP_DIR / "assets" / "Forecast360.png",
        Path.cwd() / "assets" / "Forecast360.png",
        Path(env_asst).expanduser().resolve() if env_asst else None,
    ]
    user = next((p for p in user_candidates if p and p.exists()), None)
    asst = next((p for p in asst_candidates if p and p.exists()), None)
    return user, asst

def _img_to_data_uri(path: Optional[Path]) -> Optional[str]:
    if not path or not path.exists():
        return None
    try:
        b64 = base64.b64encode(path.read_bytes()).decode("ascii")
        ext = path.suffix.lower().lstrip(".") or "png"
        mime = "image/png" if ext in ("png","apng") else ("image/jpeg" if ext in ("jpg","jpeg") else "image/svg+xml")
        return f"data:{mime};base64,{b64}"
    except Exception:
        return None

# -------------------------------------------------------------------
# Optional readers
# -------------------------------------------------------------------
try:
    from pypdf import PdfReader
    HAVE_PYPDF = True
except Exception:
    HAVE_PYPDF = False

try:
    import docx
    HAVE_DOCX = True
except Exception:
    HAVE_DOCX = False

# -------------------------------------------------------------------
# Extra optional readers (multimedia & more formats)
# -------------------------------------------------------------------
try:
    from PIL import Image
    HAVE_PIL = True
except Exception:
    HAVE_PIL = False

try:
    import pytesseract  # OCR; requires Tesseract binary installed
    HAVE_TESS = True
except Exception:
    HAVE_TESS = False

try:
    from pptx import Presentation  # PowerPoint
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

try:
    from bs4 import BeautifulSoup  # HTML parser
    HAVE_BS4 = True
except Exception:
    HAVE_BS4 = False

try:
    from striprtf.striprtf import rtf_to_text  # RTF
    HAVE_RTF = True
except Exception:
    HAVE_RTF = False

# Speech-to-text options
HAVE_WHISPER = False
try:
    from faster_whisper import WhisperModel  # recommended offline STT
    HAVE_WHISPER = True
except Exception:
    pass

try:
    import speech_recognition as sr  # fallback STT
    HAVE_SR = True
except Exception:
    HAVE_SR = False

try:
    from moviepy.editor import VideoFileClip  # video->audio
    HAVE_MOVIEPY = True
except Exception:
    HAVE_MOVIEPY = False

# -------------------------------------------------------------------
# Vector DB / Embeddings
# -------------------------------------------------------------------
import chromadb
from chromadb.config import Settings
from chromadb.utils import embedding_functions

# -------------------------------------------------------------------
# App Config / Theme
# -------------------------------------------------------------------
st.set_page_config(page_title="LLM Chat", page_icon="💬", layout="wide")
logging.getLogger("chromadb").setLevel(logging.WARNING)

# Resolve avatars once, build data URIs for CSS
USER_AVATAR_PATH, ASSIST_AVATAR_PATH = _resolve_avatar_paths()
USER_AVATAR_URI = _img_to_data_uri(USER_AVATAR_PATH)
ASSIST_AVATAR_URI = _img_to_data_uri(ASSIST_AVATAR_PATH)

# Inject CSS (white UI, light-grey sidebar, compact hero/status)
st.markdown(f"""
<style>
:root{{
  --bg:#f7f8fb;
  --sidebar-bg:#f5f7fb;              /* slightly grey sidebar */
  --hero:#ffffff;
  --panel:#ffffff;
  --text:#0b1220;
  --muted:#5d6b82;
  --accent:#2563eb;
  --border:#e7eaf2;
  --chip:#eef2ff;
  --chip-text:#334155;
  --bubble-user:#eef4ff;
  --bubble-assist:#f6f7fb;
}}
html, body, [data-testid="stAppViewContainer"]{{ background:var(--bg); color:var(--text); }}
section[data-testid="stSidebar"]{{ background:var(--sidebar-bg); border-right:1px solid var(--border); }}
section[data-testid="stSidebar"] .block-container{{ padding-top:.75rem; }}
a{{ color:var(--accent); text-decoration:none; }} a:hover{{ text-decoration:underline; }}

/* Compact the main container top padding a touch */
main .block-container{{ padding-top: .5rem; }}

/* Center the main conversation lane like ChatGPT */
.container-narrow{{ max-width:1080px; margin:0 auto; }}

/* HERO */
.hero{{
  background:var(--hero);
  border:1px solid var(--border);
  border-radius:14px;
  box-shadow:0 6px 16px rgba(16,24,40,.05);
  padding: .9rem 1rem .75rem 1rem;
  margin: .3rem auto .6rem;
}}
.hero-top{{ display:flex; align-items:center; justify-content:space-between; gap:.75rem; }}
.hero-title{{ display:flex; align-items:center; gap:.55rem; }}
.hero-title h2{{ margin:0; font-size:1.15rem;}}
.badge{{
  font-size:.72rem; color:var(--chip-text);
  background:var(--chip); border:1px solid var(--border);
  padding:.2rem .55rem; border-radius:999px;
}}
.hero-sub{{ margin-top:.3rem; color:var(--muted); font-size:.92rem; }}
.hero-body{{ margin-top:.55rem; }}

/* Nuke Streamlit's status/progress artifacts (that white pill) */
[data-testid="stStatusWidget"],
div[role="status"],
div[role="progressbar"] {{ display:none !important; }}

/* Hide Streamlit’s thin top decoration (optional) */
div[data-testid="stDecoration"]{{ display:none; }}

.status-inline{{
  width:100%; border:1px solid var(--border);
  background:#fafcff; border-radius:10px;
  padding:.5rem .7rem; font-size:.9rem; color:#111827;
  margin-top:.55rem;
}}

/* CHAT CARD */
.chat-card{{
  background:var(--panel);
  border:1px solid var(--border);
  border-radius:14px;
  box-shadow:0 6px 16px rgba(16,24,40,.05);
  overflow:hidden;
}}
.chat-scroll{{ max-height: 52vh; overflow:auto; padding:.65rem .9rem; }}
.msg{{ display:flex; align-items:flex-start; gap:.65rem; margin:.45rem 0; }}
.avatar{{
  width:32px; height:32px; border-radius:50%;
  background:#e9eef9; border:1px solid var(--border); flex:0 0 32px;
  background-size:cover; background-position:center; background-repeat:no-repeat;
}}
.avatar.user {{
  {"background-image:url('" + USER_AVATAR_URI + "');" if USER_AVATAR_URI else ""}
}}
.avatar.assistant {{
  {"background-image:url('" + ASSIST_AVATAR_URI + "');" if ASSIST_AVATAR_URI else ""}
}}
.bubble{{
  border:1px solid var(--border);
  background:var(--bubble-assist);
  padding:.8rem .95rem; border-radius:12px; max-width:860px;
  white-space:pre-wrap; line-height:1.45;
}}
.msg.user .bubble{{ background:var(--bubble-user); }}
.composer{{ padding:.6rem .75rem; border-top:1px solid var(--border); background:#fff; position:sticky; bottom:0; z-index:2; }}
</style>
""", unsafe_allow_html=True)

# -------------------------------------------------------------------
# Utilities
# -------------------------------------------------------------------
TEXT_EXTS = {".txt", ".md", ".rtf", ".html", ".htm", ".json", ".xml"}
DOC_EXTS  = {".pdf", ".docx", ".csv", ".tsv", ".xlsx", ".xlsm", ".xltx", ".pptx"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
AUDIO_EXTS = {".wav", ".mp3", ".m4a", ".flac", ".ogg", ".opus"}
VIDEO_EXTS = {".mp4", ".mov", ".mkv", ".avi", ".webm"}

SUPPORTED_EXTS = TEXT_EXTS | DOC_EXTS | IMAGE_EXTS | AUDIO_EXTS | VIDEO_EXTS

def human_time(ms: float) -> str:
    return f"{ms:.0f} ms" if ms < 1000 else f"{ms/1000:.2f} s"

def stable_hash(s: str) -> str:
    import hashlib
    return hashlib.sha256(s.encode("utf-8")).hexdigest()[:12]

def get_kb_dir() -> str:
    kb = os.path.abspath(os.path.join(".", "KB"))
    os.makedirs(kb, exist_ok=True)
    return kb

# ----- Readers -----
def read_text_file(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        try:
            with open(path, "rb") as f:
                return f.read().decode("utf-8", errors="ignore")
        except Exception:
            return ""

def read_html(path: str) -> str:
    raw = read_text_file(path)
    if not raw:
        return ""
    if HAVE_BS4:
        try:
            soup = BeautifulSoup(raw, "lxml")
            return "\n".join(s.strip() for s in soup.stripped_strings)
        except Exception:
            pass
    return re.sub(r"<[^>]+>", " ", raw)

def read_rtf(path: str) -> str:
    if not HAVE_RTF:
        return ""
    try:
        return rtf_to_text(read_text_file(path))
    except Exception:
        return ""

def read_json(path: str) -> str:
    try:
        import json
        obj = json.load(open(path, "r", encoding="utf-8", errors="ignore"))
    except Exception:
        return read_text_file(path)
    out = []
    def walk(k, v, prefix=""):
        key = f"{prefix}.{k}" if prefix else str(k)
        if isinstance(v, dict):
            for kk, vv in v.items():
                walk(kk, vv, key)
        elif isinstance(v, list):
            for i, item in enumerate(v):
                walk(str(i), item, key)
        else:
            try:
                out.append(f"{key}: {str(v)}")
            except Exception:
                pass
    if isinstance(obj, dict):
        for k, v in obj.items():
            walk(k, v)
    elif isinstance(obj, list):
        for i, v in enumerate(obj):
            walk(str(i), v)
    else:
        out.append(str(obj))
    return "\n".join(out)

def read_xml(path: str) -> str:
    try:
        from lxml import etree
        parser = etree.XMLParser(recover=True)
        tree = etree.parse(path, parser)
        return " ".join(tree.xpath("//text()"))
    except Exception:
        return re.sub(r"<[^>]+>", " ", read_text_file(path))

def read_pdf(path: str) -> str:
    if not HAVE_PYPDF: return ""
    try:
        reader = PdfReader(path)
        return "\n".join((p.extract_text() or "") for p in reader.pages)
    except Exception:
        return ""

def read_docx(path: str) -> str:
    if not HAVE_DOCX: return ""
    try:
        document = docx.Document(path)
        return "\n".join(para.text for para in document.paragraphs)
    except Exception:
        return ""

def read_pptx(path: str) -> str:
    if not HAVE_PPTX: return ""
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                except Exception:
                    continue
        return "\n".join(t.strip() for t in texts if t and t.strip())
    except Exception:
        return ""

def read_spreadsheet(path: str) -> str:
    try:
        if path.lower().endswith((".xlsx",".xlsm",".xltx")):
            df = pd.read_excel(path, engine="openpyxl")
        else:
            df = pd.read_csv(path)
        df = df.astype(str).iloc[:1000, :50]
        header = " | ".join(df.columns.tolist())
        body = "\n".join(" | ".join(row) for row in df.values.tolist())
        return f"{header}\n{body}"
    except Exception:
        return ""

def read_image_ocr(path: str) -> str:
    if not (HAVE_PIL and HAVE_TESS):
        return ""
    try:
        img = Image.open(path).convert("L")
        return pytesseract.image_to_string(img) or ""
    except Exception:
        return ""

_WHISPER_MODEL = None
def _get_whisper():
    global _WHISPER_MODEL
    if _WHISPER_MODEL or not HAVE_WHISPER:
        return _WHISPER_MODEL
    try:
        _WHISPER_MODEL = WhisperModel("base")
    except Exception:
        _WHISPER_MODEL = None
    return _WHISPER_MODEL

def read_audio_stt(path: str) -> str:
    m = _get_whisper()
    if m:
        try:
            segments, _ = m.transcribe(path, vad_filter=True)
            return " ".join(seg.text.strip() for seg in segments if getattr(seg, "text", "").strip())
        except Exception:
            pass
    if HAVE_SR:
        try:
            r = sr.Recognizer()
            with sr.AudioFile(path) as src:
                audio = r.record(src)
            try:
                return r.recognize_sphinx(audio)
            except Exception:
                return r.recognize_google(audio, show_all=False)
        except Exception:
            return ""
    return ""

def read_video_to_text(path: str) -> str:
    if not HAVE_MOVIEPY:
        return ""
    try:
        clip = VideoFileClip(path)
        tmp_wav = os.path.join(os.path.dirname(path), f"__tmp_{uuid.uuid4().hex}.wav")
        clip.audio.write_audiofile(tmp_wav, verbose=False, logger=None)
        clip.close()
        text = read_audio_stt(tmp_wav)
        try:
            os.remove(tmp_wav)
        except Exception:
            pass
        return text
    except Exception:
        return ""

def load_file_to_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()

    if ext in {".txt", ".md"}: 
        return read_text_file(path)
    if ext in {".html", ".htm"}:
        return read_html(path)
    if ext == ".rtf":
        return read_rtf(path)
    if ext == ".json":
        return read_json(path)
    if ext == ".xml":
        return read_xml(path)

    if ext == ".pdf":
        return read_pdf(path)
    if ext == ".docx":
        return read_docx(path)
    if ext in {".csv", ".tsv", ".xlsx", ".xlsm", ".xltx"}:
        return read_spreadsheet(path)
    if ext == ".pptx":
        return read_pptx(path)

    if ext in IMAGE_EXTS:
        return read_image_ocr(path)
    if ext in AUDIO_EXTS:
        return read_audio_stt(path)
    if ext in VIDEO_EXTS:
        return read_video_to_text(path)

    return ""

# ----- Chunking -----
@dataclass
class ChunkingConfig:
    chunk_size: int = 1200
    chunk_overlap: int = 200
    min_chunk_size: int = 300

def split_text_recursive(text: str, cfg: ChunkingConfig) -> List[str]:
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    if len(text) <= cfg.chunk_size: return [text]
    candidates: List[str] = []
    for delim in ["\n\n", "\n", ". "]:
        parts = text.split(delim)
        buf = ""; tmp: List[str] = []
        for p in parts:
            cand = (buf + (delim if buf else "") + p).strip()
            if len(cand) >= cfg.chunk_size:
                tmp.append(cand); buf = ""
            else:
                buf = cand
        if buf: tmp.append(buf)
        avg = sum(len(c) for c in tmp) / max(1, len(tmp))
        if avg <= cfg.chunk_size * 1.2:
            candidates = tmp; break
    if not candidates: candidates = [text]
    final_chunks: List[str] = []
    for block in candidates:
        if len(block) <= cfg.chunk_size: final_chunks.append(block)
        else:
            step = cfg.chunk_size - cfg.chunk_overlap
            for start in range(0, len(block), step):
                piece = block[start:start+cfg.chunk_size]
                if len(piece) >= cfg.min_chunk_size or (start == 0 and piece):
                    final_chunks.append(piece)
    return [c.strip() for c in final_chunks if c.strip()]

# ----- Chroma -----
def get_chroma_client(persist_dir: str) -> chromadb.Client:
    import os
    os.makedirs(persist_dir, exist_ok=True)
    return chromadb.Client(Settings(anonymized_telemetry=False, persist_directory=persist_dir))

def get_sentence_transformer_fn(model_name: str):
    return embedding_functions.SentenceTransformerEmbeddingFunction(model_name=model_name, device=None)

def ensure_collection(client: chromadb.Client, name: str, embedding_fn):
    try:
        col = client.get_collection(name=name)
        if getattr(col, "_embedding_function", None) is None:
            client.delete_collection(name)
            col = client.create_collection(name=name, embedding_function=embedding_fn, metadata={"hnsw:space":"cosine"})
    except Exception:
        col = client.create_collection(name=name, embedding_function=embedding_fn, metadata={"hnsw:space":"cosine"})
    return col

# ----- Indexing -----
def iter_files(folder: str) -> List[str]:
    paths: List[str] = []
    for ext in SUPPORTED_EXTS:
        paths.extend(glob.glob(os.path.join(folder, f"**/*{ext}"), recursive=True))
    return sorted(list(set(paths)))

def compute_kb_signature(folder: str) -> Tuple[str, int]:
    files = iter_files(folder)
    lines = []
    base = os.path.abspath(folder)
    for p in files:
        try:
            stt = os.stat(p)
            rel = os.path.relpath(os.path.abspath(p), base)
            lines.append(f"{rel}|{stt.st_size}|{int(stt.st_mtime)}")
        except Exception:
            continue
    lines.sort()
    raw = "\n".join(lines)
    return stable_hash(raw if raw else f"EMPTY-{time.time()}"), len(files)

def index_folder(
    folder: str,
    client: chromadb.Client,
    collection_name: str,
    embedding_model: str = "sentence-transformers/all-MiniLM-L6-v2",
    chunk_cfg: Optional[ChunkingConfig] = None,
    batch_size: int = 64,
) -> Tuple[int,int]:
    chunk_cfg = chunk_cfg or ChunkingConfig()
    emb_fn = get_sentence_transformer_fn(embedding_model)
    col = ensure_collection(client, collection_name, emb_fn)
    files = iter_files(folder)
    added_docs = added_chunks = 0
    ids: List[str] = []; docs: List[str] = []; metas: List[Dict] = []
    for path in files:
        text = load_file_to_text(path)
        if not text.strip(): continue
        chunks = split_text_recursive(text, chunk_cfg)
        file_id = stable_hash(path); ts = int(time.time())
        for i, ch in enumerate(chunks):
            ids.append(f"{file_id}-{i}-{ts}-{uuid.uuid4().hex[:8]}")
            docs.append(ch)
            metas.append({"source": path, "file_id": file_id, "chunk_index": i, "mtime": os.path.getmtime(path)})
            if len(ids) >= batch_size:
                col.add(ids=ids, documents=docs, metadatas=metas)
                added_chunks += len(ids); ids, docs, metas = [], [], []
        added_docs += 1
    if ids:
        col.add(ids=ids, documents=docs, metadatas=metas); added_chunks += len(ids)
    try: client.persist()
    except Exception: pass
    gc.collect(); return added_docs, added_chunks

# ----- Retrieval & Prompt -----
def retrieve_context(client: chromadb.Client, collection_name: str, query: str, k: int = 5) -> Tuple[List[str], List[Dict]]:
    col = client.get_collection(collection_name)
    out = col.query(query_texts=[query], n_results=k, include=["documents","metadatas","distances"])
    docs = out.get("documents", [[]])[0]
    metas = out.get("metadatas", [[]])[0]
    return docs, metas

def build_prompt(question: str, context_docs: List[str], meta: List[Dict], max_chars: int = 6000) -> Tuple[str,List[Tuple[str,str]]]:
    pairs: List[Tuple[str,str]] = []
    context = ""
    for d, m in zip(context_docs, meta):
        src = m.get("source", "unknown"); pairs.append((src, d))
        block = f"\n[Source: {src}]\n{d}\n"
        if len(context) + len(block) <= max_chars: context += block
        else: break
    sys_prompt = (
        "You are a careful assistant that answers ONLY using the provided context snippets. "
        "If the answer isn't in the context, say you don't know. Be concise and cite sources as [n]."
    )
    user_prompt = (
        f"Question: {question}\n\nContext:\n{context}\n"
        "Instructions:\n- Use ONLY the context above.\n- Cite sources like [1], [2] based on order."
    )
    return (sys_prompt + "\n\n" + user_prompt), pairs

# ----- LLM Backends -----
DEFAULT_OLLAMA = "llama3.2:1b-instruct-q4_K_M"
DEFAULT_CLAUDE = "claude-sonnet-4-5"

def call_ollama_chat(model: str, prompt: str, temperature: float = 0.2, timeout: int = 120) -> str:
    try:
        resp = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": model, "prompt": prompt, "temperature": temperature, "stream": False},
            timeout=timeout,
        ); resp.raise_for_status(); data = resp.json()
        return data.get("response", "").strip()
    except Exception as e:
        return f"[Ollama error] {e}"

# def call_claude(model: str, prompt: str, temperature: float = 0.2, timeout: int = 120) -> str:
#     api_key = os.getenv("ANTHROPIC_API_KEY", "")
#     if not api_key:
#         return "[Claude error] Missing ANTHROPIC_API_KEY environment variable."
#     try:
#         from anthropic import Anthropic
#         client = Anthropic(api_key=api_key, timeout=timeout)
#         msg = client.messages.create(
#             model=model,
#             max_tokens=800,
#             temperature=temperature,
#             system="You are a concise, citation-focused assistant.",
#             messages=[{"role": "user", "content": prompt}],
#         )
#         parts: List[str] = []
#         for block in msg.content:
#             t = getattr(block, "text", None) or (block.get("text") if isinstance(block, dict) else None)
#             if t: parts.append(t)
#         return "\n".join(parts).strip()
#     except Exception as e:
#         return f"[Claude error] {e}"

import os
from typing import List, Optional

def _get_streamlit_secret(name: str) -> Optional[str]:
    """
    Helper: return secret from streamlit.secrets if Streamlit is available,
    otherwise return None.
    """
    try:
        import streamlit as _st
        # streamlit.secrets behaves like a dict; use .get to avoid KeyError
        return _st.secrets.get(name)
    except Exception:
        return None


def _extract_text_from_msg(msg) -> str:
    """
    Try a few heuristics to extract readable text from an Anthropic response.
    Handles object-like responses with .content, dict responses, lists of blocks,
    and plain strings.
    """
    # 1) object with attribute .content
    content = None
    if hasattr(msg, "content"):
        try:
            content = msg.content
        except Exception:
            content = None

    # 2) dict-like response
    if content is None and isinstance(msg, dict):
        # common keys: "completion", "output", "text", "content"
        for k in ("completion", "output", "text", "content"):
            if k in msg:
                content = msg[k]
                break

    # 3) fallback to msg itself (string or iterable)
    if content is None:
        content = msg

    # Now normalize content into a string
    parts: List[str] = []

    # If it's a plain string, just return it
    if isinstance(content, str):
        return content.strip()

    # If it's an iterable (list/tuple), iterate blocks
    if isinstance(content, (list, tuple)):
        for block in content:
            if isinstance(block, str):
                text = block
            elif isinstance(block, dict):
                # try several keys commonly used for text
                text = (
                    block.get("text")
                    or block.get("content")
                    or block.get("output")
                    or block.get("data")
                    or None
                )
            else:
                # object-like block
                text = getattr(block, "text", None) or getattr(block, "content", None)
            if text:
                parts.append(str(text).strip())
        return "\n".join(parts).strip()

    # If it's a dict-like content now, try to pull a text field
    if isinstance(content, dict):
        for k in ("text", "content", "output", "completion", "data"):
            if k in content:
                val = content[k]
                if isinstance(val, str):
                    return val.strip()
                if isinstance(val, (list, tuple)):
                    return "\n".join(map(str, val)).strip()
        # last resort: return stringified dict
        return str(content).strip()

    # Fallback: string-convert whatever it is
    try:
        return str(content).strip()
    except Exception:
        return ""


def call_claude(
    model: str,
    prompt: str,
    temperature: float = 0.2,
    timeout: int = 120,
    max_tokens: int = 800,
) -> str:
    """
    Call Anthropic/Claude and return a text response. Looks for ANTHROPIC_API_KEY
    in the environment or streamlit.secrets["ANTHROPIC_API_KEY"].

    Returns a readable string on success, or an error message beginning with
    "[Claude error]" on failure.
    """
    # 1) find the API key (env first, then Streamlit secrets)
    api_key = os.environ.get("ANTHROPIC_API_KEY") or _get_streamlit_secret("ANTHROPIC_API_KEY")
    if not api_key:
        return "[Claude error] Missing ANTHROPIC_API_KEY environment variable or streamlit secret."

    try:
        # Import lazily so the module isn't required unless this function runs
        from anthropic import Anthropic

        # Construct client (some anthopic versions may accept timeout kw)
        client = Anthropic(api_key=api_key, timeout=timeout)

        # Preferred call pattern (keeps your previous structure)
        # If your anthopic SDK version uses a slightly different API, you can
        # substitute the appropriate client method here.
        msg = client.messages.create(
            model=model,
            max_tokens=max_tokens,
            temperature=temperature,
            system="You are a concise, citation-focused assistant.",
            messages=[{"role": "user", "content": prompt}],
        )

        # Normalise/parse response text robustly
        return _extract_text_from_msg(msg) or "[Claude error] Empty response."

    except Exception as e:
        # Provide a short, non-sensitive message (avoid echoing secrets)
        return f"[Claude error] {type(e).__name__}: {str(e)}"
# -------------------------------------------------------------------
# Settings defaults
# -------------------------------------------------------------------
def settings_defaults():
    kb_dir = get_kb_dir()
    return {
        "persist_dir": ".chroma",
        "collection_name": f"kb-{stable_hash(kb_dir)}",
        "base_folder": kb_dir,
        "emb_model": "sentence-transformers/all-MiniLM-L6-v2",
        "chunk_cfg": ChunkingConfig(chunk_size=1200, chunk_overlap=200, min_chunk_size=300),
        "backend": "Claude (Anthropic)",
        "ollama_model": DEFAULT_OLLAMA,
        "claude_model": DEFAULT_CLAUDE,
        "temperature": 0.2,
        "top_k": 5,
        "auto_index_min_interval_sec": 8,
    }

# -------------------------------------------------------------------
# Auto-index (renders status inline on RHS) — NO st.status (no white pill)
# -------------------------------------------------------------------
def auto_index_if_needed(status_placeholder=None):
    folder = st.session_state["base_folder"]
    persist = st.session_state["persist_dir"]
    colname = st.session_state["collection_name"]
    emb_model = st.session_state["emb_model"]
    min_gap = int(st.session_state.get("auto_index_min_interval_sec", 8))

    client = get_chroma_client(persist)

    sig_now, file_count = compute_kb_signature(folder)
    last_sig = st.session_state.get("_kb_last_sig")
    last_time = float(st.session_state.get("_kb_last_index_ts", 0.0))
    now = time.time()

    need_index = (last_sig != sig_now) or (last_sig is None)
    throttled = (now - last_time) < min_gap

    target = status_placeholder if status_placeholder is not None else st

    if need_index and not throttled:
        try:
            n_docs, n_chunks = index_folder(
                folder=folder, client=client, collection_name=colname,
                embedding_model=emb_model, chunk_cfg=st.session_state["chunk_cfg"],
            )
    st.session_state["_kb_last_sig"] = sig_now
    st.session_state["_kb_last_index_ts"] = now
    st.session_state["_kb_last_counts"] = {"files": file_count, "docs": n_docs, "chunks": n_chunks}
            label = f"Indexed: <b>{n_docs}</b> files processed, <b>{n_chunks}</b> chunks"
        except Exception as e:
            label = f"Auto-index failed: <b>{e}</b>"
        target.markdown(f'<div class="status-inline">{label}</div>', unsafe_allow_html=True)
    else:
        ts = st.session_state.get("_kb_last_index_ts")
        when = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(ts)) if ts else "—"
        target.markdown(
            f'<div class="status-inline">Auto-index is <b>ON</b> · Files: <b>{file_count}</b> · Last indexed: <b>{when}</b> · Collection: <code>{colname}</code></div>',
            unsafe_allow_html=True
        )
    return client

# -------------------------------------------------------------------
# Main
# -------------------------------------------------------------------
def main():
    sdef = settings_defaults()
    for k, v in sdef.items():
        st.session_state.setdefault(k, v)

    if st.session_state.pop("_clear_compose", False):
        st.session_state["compose_input"] = ""

    st.session_state.setdefault("messages", [
        {"role": "assistant", "content": "Hi! I’m your knowledge assistant. Ask anything about your Knowledge Base."}
    ])

    # Sidebar
    with st.sidebar:
        logo_path = _resolve_logo_path()
if logo_path:
    try:
        st_image_compat(
            str(logo_path),
            caption="iSOFT ANZ Pvt Ltd",
            use_column_width=True,
        )
    except Exception:
        # Fallback to plain st.image if compatibility wrapper fails
        try:
            st.image(str(logo_path), caption="iSOFT ANZ Pvt Ltd")
        except Exception:
            pass

# Chroma admin helper (best-effort, don't let it crash the sidebar)
try:
    show_chroma_admin_ui()
except Exception:
    pass

st.subheader("⚙️ Settings")
st.caption("Auto-index is enabled. Edit paths/models below if needed.")
st.session_state["base_folder"] = st.text_input("Knowledge Base", value=st.session_state.get("base_folder", get_kb_dir()))
st.session_state["persist_dir"] = st.text_input("Chroma persist", value=st.session_state["persist_dir"])
        st.session_state["collection_name"] = st.text_input("Collection", value=st.session_state["collection_name"])
        st.session_state["emb_model"] = st.selectbox(
            "Embedding model",
            [
                "sentence-transformers/all-MiniLM-L6-v2",
                "sentence-transformers/all-MiniLM-L12-v2",
                "sentence-transformers/multi-qa-MiniLM-L6-cos-v1",
            ],
            index=0,
        )
        st.session_state["backend"] = st.radio("LLM", ["Claude (Anthropic)", "Ollama (local)"], index=0)
        if st.session_state["backend"].startswith("Claude"):
            st.session_state["claude_model"] = st.text_input("Claude model", value=st.session_state["claude_model"])
        else:
            st.session_state["ollama_model"] = st.text_input("Ollama model", value=st.session_state["ollama_model"])
        st.session_state["temperature"] = st.slider("Temperature", 0.0, 1.0, float(st.session_state["temperature"]), 0.05)
        st.session_state["top_k"] = st.slider("Top-K", 1, 15, int(st.session_state["top_k"]))
        st.session_state["auto_index_min_interval_sec"] = st.number_input(
            "Auto-index min interval (sec)", min_value=1, max_value=300,
            value=int(st.session_state["auto_index_min_interval_sec"]), step=1
        )
        st.caption("Tip: Drop PDFs, DOCX, CSV/XLSX, text files, images, audio, videos, PPTX, HTML, RTF, JSON, XML into ./KB — indexing is automatic (if dependencies are installed).")

    # ------------------- Centered main lane -------------------
    # st.markdown('<div class="container-narrow">', unsafe_allow_html=True)

    # HERO (chat lives inside)
    # st.markdown('<div class="hero">', unsafe_allow_html=True)
    st.markdown(
        """
        <div class="hero-top">
          <div class="hero-title">
            <span style="font-size:1.05rem">💬</span>
            <h2>Chat with LLM</h2>
            <span class="badge">RAG • Chroma • Streamlit</span>
          </div>
        </div>
        <div class="hero-sub">Ask questions grounded in your Knowledge Base. We’ll search your indexed files and cite context.</div>
        """,
        unsafe_allow_html=True
    )
    hero_status = st.container()
    client = auto_index_if_needed(status_placeholder=hero_status)

    # --- Chat inside hero ---
    st.markdown('<div class="hero-body">', unsafe_allow_html=True)

    st.markdown('<div class="chat-card">', unsafe_allow_html=True)
    st.markdown('<div class="chat-scroll">', unsafe_allow_html=True)
    for m in st.session_state["messages"]:
        role = m["role"]
        who = "user" if role == "user" else "assistant"
        st.markdown(
            f'<div class="msg {who}"><div class="avatar {who}"></div><div class="bubble">{m["content"]}</div></div>',
            unsafe_allow_html=True
        )
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="composer">', unsafe_allow_html=True)
    col_input, col_send = st.columns([1, 0.22])
    with col_input:
        compose_val = st.session_state.get("compose_input", "")
        _ = st.text_area("Message", key="compose_input", value=compose_val,
                         placeholder="Type your question…", label_visibility="collapsed")
    with col_send:
        send_clicked = st.button("Send", use_container_width=True)
    st.markdown('</div>', unsafe_allow_html=True)  # /composer
    st.markdown('</div>', unsafe_allow_html=True)  # /chat-card
    st.markdown('</div>', unsafe_allow_html=True)  # /hero-body
    st.markdown('</div>', unsafe_allow_html=True)  # /hero

    # ------------------- End centered main lane -------------------
    # st.markdown('</div>', unsafe_allow_html=True)

    # SEND -> RAG -> LLM
    if send_clicked and (text := st.session_state.get("compose_input", " ").strip()):
        st.session_state["messages"].append({"role": "user", "content": text})
        try:
            docs, metas = retrieve_context(client, st.session_state["collection_name"], text, k=int(st.session_state["top_k"]))
        except Exception as e:
            st.session_state["messages"].append({"role": "assistant", "content": f"Retrieval failed: {e}"})
            st.session_state["_clear_compose"] = True
            st.rerun()

        if not docs:
            st.session_state["messages"].append({"role": "assistant", "content": "I couldn’t find context yet. Add files to ./KB (or your chosen folder) and give it a moment for auto-index."})
            st.session_state["_clear_compose"] = True
            st.rerun()

        prompt, _pairs = build_prompt(text, docs, metas)
        t0 = time.time()
        if st.session_state["backend"].startswith("Ollama"):
            out = call_ollama_chat(st.session_state["ollama_model"], prompt, temperature=float(st.session_state["temperature"]))
        else:
            out = call_claude(st.session_state["claude_model"], prompt, temperature=float(st.session_state["temperature"]))
        t1 = time.time()
        answer = f"{out}\n\n_(Answered in {human_time((t1-t0)*1000)} using {st.session_state['backend']})_"
        st.session_state["messages"].append({"role": "assistant", "content": answer})
        st.session_state["_clear_compose"] = True
        st.rerun()

if __name__ == "__main__":
    main()
